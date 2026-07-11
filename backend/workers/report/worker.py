import uuid
import uuid6
from datetime import datetime, UTC
from pathlib import Path
import tempfile
import os

from celery import shared_task
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import seaborn as sns
import io

from backend.common.storage import minio_client
from backend.common.storage.minio_client import upload_file
from backend.common.schemas.worker_outputs import ReportWorkerOutputs
from backend.common.schemas.worker_envelope import WorkerEnvelope, WorkerSource
from backend.common.schemas.canonical_result import CanonicalResult



# chart generation
def generate_charts(result: CanonicalResult) -> dict:
    # Generates in-memory chart images
    charts = {}

    # 1. Trial Status Distribution
    statuses = [t.status for t in result.trials]
    if statuses:
        plt.figure(figsize=(6, 6))
        # Count statuses
        status_counts = {}
        for s in statuses:
            status_counts[s] = status_counts.get(s, 0) + 1
        
        plt.pie(status_counts.values(), labels=status_counts.keys(), autopct='%1.1f%%', colors=sns.color_palette('pastel'))
        plt.title(f'Trial Status Distribution ({len(statuses)} Trials)')
        
        buf = io.BytesIO()
        plt.savefig(buf, format='png')
        buf.seek(0)
        charts['status_pie'] = buf
        plt.close()

    # 2. Phase Distribution
    phases = [t.phase for t in result.trials]
    if phases:
        plt.figure(figsize=(8, 5))
        phase_counts = {}
        for p in phases:
            phase_counts[p] = phase_counts.get(p, 0) + 1
            
        sns.barplot(x=list(phase_counts.keys()), y=list(phase_counts.values()), palette='viridis')
        plt.title('Trials by Phase')
        plt.ylabel('Count')
        
        buf = io.BytesIO()
        plt.savefig(buf, format='png')
        buf.seek(0)
        charts['phase_bar'] = buf
        plt.close()

    # 3. Top Conditions
    conditions = [t.condition for t in result.trials]
    if conditions:
        plt.figure(figsize=(10, 6))
        cond_counts = {}
        for c in conditions:
            cond_counts[c] = cond_counts.get(c, 0) + 1
        
        # Sort and take top 5
        sorted_conds = sorted(cond_counts.items(), key=lambda x: x[1], reverse=True)[:5]
        
        sns.barplot(x=[x[1] for x in sorted_conds], y=[x[0] for x in sorted_conds], palette='rocket')
        plt.title('Top 5 Conditions Studied')
        plt.xlabel('Count')
        
        buf = io.BytesIO()
        plt.savefig(buf, format='png')
        buf.seek(0)
        charts['condition_bar'] = buf
        plt.close()

    return charts

# pdf generation
def draw_wrapped_text(c, text, x, y, max_width, line_height=14):
    """
    Helper to draw text wrapped within max_width.
    Returns the new y position after drawing.
    """
    from reportlab.lib.utils import simpleSplit
    # reportlab simpleSplit(text, fontName, fontSize, maxWidth)
    lines = simpleSplit(text, c._fontname, c._fontsize, max_width)
    for line in lines:
        if y < 50: # Simple page break check (not fully robust for multi-page but safe for proto)
            c.showPage()
            y = 750
            c.setFont("Helvetica", 12)
        c.drawString(x, y, line)
        y -= line_height
    return y

def generate_pdf(report_path: Path, result: CanonicalResult):
    
    charts = generate_charts(result) # Generate charts first

    c = canvas.Canvas(
        str(report_path),
        pagesize=letter
    )
    width, height = letter
    max_width = width - 100 # 50px margin on each side

    y = 750
    c.setFont("Helvetica-Bold", 18)
    c.drawString(50, y, f"Pharma-Agent Comprehensive Report: {result.molecule}")
    y -= 30
    c.setFont("Helvetica-Oblique", 10)
    c.drawString(50, y, f"Generated on {datetime.now().strftime('%Y-%m-%d')}")
    y -= 40

    # ---------------------------------------------------------
    # SECTION 1: CLINICAL INTELLIGENCE (LIVE)
    # ---------------------------------------------------------
    c.setFont("Helvetica-Bold", 14)
    c.setFillColorRGB(0, 0, 0.5) # Dark Blue
    c.drawString(50, y, "1. CLINICAL INTELLIGENCE (Active)")
    c.setFillColorRGB(0, 0, 0) # Reset
    y -= 25
    c.setFont("Helvetica", 10)
    c.drawString(50, y, "Source: Clinical Trials Worker & Synthesis Engine")
    y -= 20

    # Summary
    c.setFont("Helvetica-Bold", 12)
    c.drawString(50, y, "Executive Summary:")
    y -= 20
    c.setFont("Helvetica", 10)
    y = draw_wrapped_text(c, result.trial_summary, 50, y, max_width)
    y -= 20

    # Charts
    if y < 300: c.showPage(); y = 750
    
    if 'status_pie' in charts:
        c.drawImage(report_lab_image_reader(charts['status_pie']), 50, y - 180, width=180, height=180)
    if 'phase_bar' in charts:
        c.drawImage(report_lab_image_reader(charts['phase_bar']), 250, y - 180, width=250, height=160)
    
    if charts: y -= 200

    if 'condition_bar' in charts:
        if y < 200: c.showPage(); y = 750
        c.drawImage(report_lab_image_reader(charts['condition_bar']), 100, y - 180, width=350, height=180)
        y -= 200
    
    # Key Findings
    if y < 100: c.showPage(); y = 750
    c.setFont("Helvetica-Bold", 12)
    c.drawString(50, y, "Key Clinical Findings:")
    y -= 20
    c.setFont("Helvetica", 10)
    for k in result.key_findings:
        y = draw_wrapped_text(c, f"• {k}", 50, y, max_width, line_height=12)
    y -= 30

    
    y -= 20

    # ---------------------------------------------------------
    # SECTION 2: MARKET INTELLIGENCE (Active)
    # ---------------------------------------------------------
    if y < 150: c.showPage(); y = 750
    c.setFont("Helvetica-Bold", 14)
    c.setFillColorRGB(0, 0, 0.5)
    c.drawString(50, y, "2. MARKET INTELLIGENCE (Active)")
    c.setFillColorRGB(0, 0, 0)
    y -= 25
    c.setFont("Helvetica", 10)
    c.drawString(50, y, f"Market Size: {result.market_data.get('market_size', 'N/A') if result.market_data else 'N/A'}")
    y -= 20
    c.drawString(50, y, f"Patent Status: {result.market_data.get('patent_status', 'N/A') if result.market_data else 'N/A'}")
    y -= 25
    
    c.setFont("Helvetica-Bold", 12)
    c.drawString(50, y, "Competitor Landscape:")
    y -= 20
    c.setFont("Helvetica", 10)
    market_findings = result.market_data.get('key_findings', []) if result.market_data else []
    for finding in market_findings:
        y = draw_wrapped_text(c, f"• {finding}", 50, y, max_width, line_height=12)
    y -= 30

    # ---------------------------------------------------------
    # SECTION 5: INTERNAL DATA CORRELATION
    # ---------------------------------------------------------
    if y < 150: c.showPage(); y = 750
    c.setFont("Helvetica-Bold", 14)
    c.setFillColorRGB(0, 0, 0.5)
    c.drawString(50, y, "5. INTERNAL DATA CORRELATION")
    c.setFillColorRGB(0, 0, 0)
    y -= 25
    c.setFont("Helvetica-Oblique", 10)
    c.drawString(50, y, "Deployment Pending: Internal Summarizer (Proprietary Data)")
    y -= 20

    c.setDash(1, 2)
    c.rect(50, y-60, max_width, 60)
    c.drawString(70, y-30, "Subject: Internal Lab Results & Historical Matches")
    c.drawString(70, y-45, "Status: Waiting for Internal Summarizer module...")
    c.setDash()
    y -= 90

    c.showPage()
    c.save() 

def report_lab_image_reader(buf):
    from reportlab.lib.utils import ImageReader
    return ImageReader(buf) 

# ppt generation
def generate_ppt(ppt_path: Path, result: CanonicalResult):

    prs= Presentation()

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"Evidence Summary: {result.molecule}"
    subtitle.text = "Generated by Pharma-Agent"

    # Slide 2: Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    body = slide.shapes.placeholders[1].text_frame
    body.text = "Overview"
    body.add_paragraph().text = f"Molecule: {result.molecule}"
    body.add_paragraph().text = f"Data Completeness Score: {result.data_completeness_score}"
    body.add_paragraph().text = f"Confidence Overall: {result.confidence_overall}"

    # Slide 3: Trial Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    body = slide.shapes.placeholders[1].text_frame
    body.text = "Clinical Trial Summary"
    for line in result.trial_summary.split("\n"):
        body.add_paragraph().text = line

    # Slide 4: Key Findings
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    body = slide.shapes.placeholders[1].text_frame
    body.text = "Key Findings"
    for k in result.key_findings:
        body.add_paragraph().text = f"• {k}"

    # Slide 5: Recommendations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    body = slide.shapes.placeholders[1].text_frame
    body.text = "Suggested Follow-Up"
    for s in result.suggested_follow_up:
        body.add_paragraph().text = f"• {s}"

    # Charts Slide
    charts = generate_charts(result)
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Blank
    title = slide.shapes.title
    title.text = "Trial Landscape"
    
    if 'status_pie' in charts:
        slide.shapes.add_picture(charts['status_pie'], Inches(0.5), Inches(2), width=Inches(4))
    if 'phase_bar' in charts:
        slide.shapes.add_picture(charts['phase_bar'], Inches(5), Inches(2), width=Inches(4))

    if 'condition_bar' in charts:
        slide = prs.slides.add_slide(prs.slide_layouts[5]) 
        title = slide.shapes.title
        title.text = "Conditions Landscape"
        slide.shapes.add_picture(charts['condition_bar'], Inches(2), Inches(2), width=Inches(6))

    # SWOT Analysis Slide
    if result.swot_analysis:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "SWOT Analysis"
        body = slide.shapes.placeholders[1].text_frame
        for cat, items in result.swot_analysis.items():
            p = body.add_paragraph()
            p.text = cat.upper()
            p.font.bold = True
            for item in items[:2]:
                body.add_paragraph().text = f"- {item}"

    # Risk Assessment Slide
    if result.risk_assessment:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Risk Assessment"
        body = slide.shapes.placeholders[1].text_frame
        body.text = result.risk_assessment

    # Slide: Key Findings
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    body = slide.shapes.placeholders[1].text_frame
    body.text = "Key Findings"
    for k in result.key_findings:
        body.add_paragraph().text = f"• {k}"

    prs.save(str(ppt_path))

# celery worker
@shared_task(name="workers.report.worker.run")
def run_report_worker(job_id: str, task_id: str, params: dict):
    # generate pdf and ppt using CanonicalResult 
    # upload artifacts to minio and 
    # return WorkerEnvelope with artifact metadata

    job_uuid = uuid.UUID(job_id)
    task_uuid = uuid.UUID(task_id)

    # 1. load CanonicalResult
    result = CanonicalResult.model_validate(params["canonical_result"])

    # 2. generate files locally
    pdf_name = f"{job_id}_report.pdf"
    ppt_name = f"{job_id}_slides.pptx"

    temp_dir = tempfile.gettempdir()
    pdf_path = Path(temp_dir) / pdf_name
    ppt_path = Path(temp_dir) / ppt_name

    generate_pdf(pdf_path, result)
    generate_ppt(ppt_path, result)

    # 3. Upload to MinIO — gracefully fall back to local paths if storage is full
    try:
        pdf_upload = upload_file(
            bucket="artifacts",
            object_name=f"{job_id}_report.pdf",
            file_path=str(pdf_path)
        )
        pdf_uri = pdf_upload["uri"]
    except Exception as e:
        print(f"[ReportWorker] WARNING: MinIO upload failed for PDF ({e}). Using local path.")
        pdf_uri = str(pdf_path)

    try:
        ppt_upload = upload_file(
            bucket="artifacts",
            object_name=f"{job_id}_slides.pptx",
            file_path=str(ppt_path)
        )
        ppt_uri = ppt_upload["uri"]
    except Exception as e:
        print(f"[ReportWorker] WARNING: MinIO upload failed for PPT ({e}). Using local path.")
        ppt_uri = str(ppt_path)

    # 4. worker output schema
    outputs = ReportWorkerOutputs(
        pdf_uri=pdf_uri,
        ppt_uri=ppt_uri
    )

    # 5. return workerenvelope
    envelope = WorkerEnvelope(
        job_id= job_uuid,
        task_id= task_uuid,
        worker="report",
        status="ok",
        confidence=1.0,
        timestamp=datetime.now(UTC),
        outputs=outputs.model_dump(),
        sources=[
            WorkerSource(
                type="generated",
                title="Report Worker Generated Files",
                uri=pdf_uri,
                retrieved_at=datetime.now(UTC)
            )
        ],
        notes=None
    )

    return envelope.model_dump(mode='json')